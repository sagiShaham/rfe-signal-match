import os, io, csv, json, uuid, asyncio, sqlite3, itertools, re, threading
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, List
from concurrent.futures import ThreadPoolExecutor
from difflib import SequenceMatcher
from collections import defaultdict

from fastapi import FastAPI, HTTPException, BackgroundTasks, UploadFile, File, Form
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, HTMLResponse
from pydantic import BaseModel
from dotenv import load_dotenv

load_dotenv()

app = FastAPI(title="RFE Signal Match")
DB_PATH = "rfe_dedup.db"
jobs: Dict[str, Dict] = {}
executor = ThreadPoolExecutor(max_workers=4)

DEFAULT_WEIGHTS = {
    "weight_release_notes": 1.0,
    "weight_ado_current_pi": 0.9,
    "weight_ado_backlog": 0.6,
    "weight_confluence_dated": 0.5,
    "weight_roadmap": 0.4,
    "weight_confluence_undated": 0.3,
}

# ─── Database ─────────────────────────────────────────────────────────────────

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS rfe_pulls (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            run_id TEXT NOT NULL,
            case_number TEXT,
            subject TEXT,
            description TEXT,
            account_name TEXT,
            account_arr REAL,
            status TEXT,
            domain TEXT,
            sub_domain TEXT,
            severity TEXT,
            created_date TEXT,
            pulled_at TEXT
        );
        CREATE TABLE IF NOT EXISTS rfe_clusters (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            run_id TEXT NOT NULL,
            cluster_title TEXT,
            confidence_score REAL,
            confidence_band TEXT,
            member_case_numbers TEXT,
            total_arr REAL,
            account_count INTEGER,
            status TEXT DEFAULT 'pending',
            match_status TEXT DEFAULT 'no_match',
            match_reason TEXT,
            score_breakdown TEXT,
            created_at TEXT
        );
        CREATE TABLE IF NOT EXISTS email_log (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            run_id TEXT, cluster_id INTEGER, case_number TEXT,
            account_name TEXT, recipient_email TEXT, subject TEXT,
            body TEXT, status TEXT, sent_at TEXT
        );
        CREATE TABLE IF NOT EXISTS run_meta (
            run_id TEXT PRIMARY KEY,
            started_at TEXT, completed_at TEXT,
            rfe_count INTEGER DEFAULT 0, cluster_count INTEGER DEFAULT 0,
            emails_sent INTEGER DEFAULT 0,
            status TEXT DEFAULT 'running', source TEXT DEFAULT 'csv'
        );
        CREATE TABLE IF NOT EXISTS app_config (key TEXT PRIMARY KEY, value TEXT);
        CREATE TABLE IF NOT EXISTS context_docs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            doc_type TEXT, title TEXT, content TEXT,
            source_url TEXT, imported_at TEXT, doc_date TEXT
        );
    """)
    conn.commit(); conn.close()

def get_cfg(key: str, default: str) -> str:
    conn = get_db(); cur = conn.cursor()
    cur.execute("SELECT value FROM app_config WHERE key=?", (key,))
    row = cur.fetchone(); conn.close()
    return row["value"] if row else os.getenv(key, default)

def set_cfg(key: str, value: str):
    conn = get_db()
    conn.execute("INSERT OR REPLACE INTO app_config (key,value) VALUES (?,?)", (key, value))
    conn.commit(); conn.close()

def migrate_db():
    conn = get_db()
    for sql in [
        "ALTER TABLE run_meta ADD COLUMN source TEXT DEFAULT 'csv'",
        "ALTER TABLE run_meta ADD COLUMN filtered_rfe_count INTEGER DEFAULT 0",
        "ALTER TABLE rfe_clusters ADD COLUMN match_status TEXT DEFAULT 'no_match'",
        "ALTER TABLE rfe_clusters ADD COLUMN match_reason TEXT",
        "ALTER TABLE rfe_clusters ADD COLUMN score_breakdown TEXT",
        "ALTER TABLE rfe_clusters ADD COLUMN confidence_sentence TEXT",
        "ALTER TABLE context_docs ADD COLUMN doc_date TEXT",
    ]:
        try: conn.execute(sql)
        except sqlite3.OperationalError: pass
    conn.commit(); conn.close()

init_db(); migrate_db()

# ─── CSV column detection ─────────────────────────────────────────────────────

SF_COLUMN_ALIASES = {
    "case_number":  ["case number","casenumber","case_number","case #","case#","case id"],
    "subject":      ["subject","case subject","rfe subject","title","case title"],
    "description":  ["description","case description","desc","details","body","content"],
    "account_name": ["account name","account: account name","account: name","accountname","company","customer"],
    "account_arr":  ["account's arr","accounts arr","annual revenue","arr","annualrevenue","account arr"],
    "status":       ["status","case status"],
    "domain":       ["product domain","product_domain__c","domain"],
    "sub_domain":   ["sub-domain","sub domain","subdomain","sub_domain__c"],
    "severity":     ["severity","case severity","priority"],
    "created_date": ["date/time opened","created date","createddate","opened date","open date","created"],
}

def detect_columns(headers: List[str]) -> Dict[str, Optional[str]]:
    normalized = {h: h.lower().strip() for h in headers}
    mapping: Dict[str, Optional[str]] = {}
    for field, aliases in SF_COLUMN_ALIASES.items():
        found = None
        for h, h_low in normalized.items():
            if any(a == h_low for a in aliases):
                found = h; break
        if not found:
            for h, h_low in normalized.items():
                if any(a in h_low for a in aliases):
                    found = h; break
        mapping[field] = found
    return mapping

def parse_arr(val: str) -> Optional[float]:
    if not val: return None
    cleaned = re.sub(r"[^\d.]", "", str(val))
    try: return float(cleaned)
    except ValueError: return None

# ─── Text similarity ──────────────────────────────────────────────────────────

STOP_WORDS = {
    "the","a","an","and","or","but","in","on","at","to","for","of","with",
    "is","are","was","were","be","been","has","have","had","do","does","did",
    "will","would","could","should","may","might","can","this","that","these",
    "those","it","its","we","our","they","their","as","by","from","not","no",
    "request","feature","ability","option","support","add","allow","enable",
    "provide","need","want","customer","user","please","would like"
}

def tokenize(text: str) -> set:
    words = re.findall(r'\b[a-z]{3,}\b', (text or "").lower())
    return {w for w in words if w not in STOP_WORDS}

def basic_similarity(rfe1: dict, rfe2: dict) -> dict:
    # Use subject only — more focused than full description for clustering
    t1 = (rfe1.get('subject') or '').strip()
    t2 = (rfe2.get('subject') or '').strip()
    w1, w2 = tokenize(t1), tokenize(t2)
    if not w1 or not w2: return {"score": 0.0, "reason": "Insufficient text.", "common": []}
    jaccard = len(w1 & w2) / len(w1 | w2)
    seq = SequenceMatcher(None, t1.lower(), t2.lower()).ratio()
    score = round(jaccard * 0.65 + seq * 0.35, 3)
    common = sorted(w1 & w2)[:7]
    reason = f"Shared key terms: {', '.join(common)}." if common else "Low lexical overlap."
    return {"score": score, "reason": reason, "common": common}

async def claude_similarity(client, rfe1: dict, rfe2: dict) -> dict:
    import anthropic as _anthropic
    prompt = (
        "You are a product deduplication engine. Given two RFE descriptions, return a JSON object with:\n"
        "- score: float 0-1\n- reason: one sentence\n- common: list of up to 5 key shared terms\n"
        "Only compare the core product ask, not writing style.\n"
        f"RFE 1: {rfe1.get('subject','')} — {(rfe1.get('description') or '')[:600]}\n"
        f"RFE 2: {rfe2.get('subject','')} — {(rfe2.get('description') or '')[:600]}\n\n"
        "Return only valid JSON."
    )
    response = await client.messages.create(
        model="claude-opus-4-5", max_tokens=200,
        messages=[{"role": "user", "content": prompt}],
    )
    text = response.content[0].text.strip()
    if "```" in text:
        text = text.split("```")[1]
        if text.startswith("json"): text = text[4:]
    return json.loads(text.strip())

# ─── Union-Find ───────────────────────────────────────────────────────────────

class UnionFind:
    def __init__(self, n):
        self.parent = list(range(n)); self.rank = [0] * n
    def find(self, x):
        while self.parent[x] != x:
            self.parent[x] = self.parent[self.parent[x]]; x = self.parent[x]
        return x
    def union(self, x, y):
        rx, ry = self.find(x), self.find(y)
        if rx == ry: return
        if self.rank[rx] < self.rank[ry]: rx, ry = ry, rx
        self.parent[ry] = rx
        if self.rank[rx] == self.rank[ry]: self.rank[rx] += 1

# ─── Scoring config ───────────────────────────────────────────────────────────

def get_scoring_config() -> dict:
    cfg = {}
    for key, default in DEFAULT_WEIGHTS.items():
        cfg[key] = float(get_cfg(key, str(default)))
    cfg['recency_decay'] = get_cfg('recency_decay', 'true').lower() == 'true'
    cfg['threshold_high'] = float(get_cfg('threshold_high', '0.75'))
    cfg['threshold_medium'] = float(get_cfg('threshold_medium', '0.45'))
    cfg['similarity_threshold'] = float(get_cfg('SIMILARITY_THRESHOLD', '0.55'))
    cfg['high_confidence_threshold'] = float(get_cfg('HIGH_CONFIDENCE_THRESHOLD', '0.85'))
    # Match status thresholds — minimum weighted score to claim each status
    cfg['match_thresh_delivered']   = float(get_cfg('match_thresh_delivered',   '0.45'))
    cfg['match_thresh_in_pi']       = float(get_cfg('match_thresh_in_pi',       '0.35'))
    cfg['match_thresh_planned']     = float(get_cfg('match_thresh_planned',     '0.28'))
    return cfg

# ─── Match status determination ───────────────────────────────────────────────

def determine_match_status(cluster_text: str, context_docs: list, cfg: dict) -> tuple:
    """Returns (status, reason, breakdown_dict)"""
    tokens = tokenize(cluster_text)
    if not tokens:
        return 'no_match', 'Insufficient text to match', {}

    t_del     = cfg.get('match_thresh_delivered', 0.45)
    t_pi      = cfg.get('match_thresh_in_pi',     0.35)
    t_planned = cfg.get('match_thresh_planned',   0.28)
    PRIORITY = [
        ('release_notes',  'delivered',     'weight_release_notes',      t_del),
        ('ado_current_pi', 'in_current_pi', 'weight_ado_current_pi',     t_pi),
        ('ado_backlog',    'planned',        'weight_ado_backlog',        t_planned),
        ('roadmap',        'planned',        'weight_roadmap',            t_planned),
        ('confluence',     'planned',        'weight_confluence_undated', t_planned),
    ]
    breakdown = {}

    for doc_type, status_label, weight_key, min_score in PRIORITY:
        type_docs = [d for d in context_docs if d['doc_type'] == doc_type]
        base_weight = cfg.get(weight_key, 0.3)
        best_score, best_doc = 0, None

        for doc in type_docs:
            content = (doc.get('content') or '')[:20000]
            doc_tokens = tokenize(content)
            if not doc_tokens: continue
            # Use containment: fraction of RFE terms that appear in the doc.
            # Jaccard (|A∩B|/|A∪B|) is wrong here — a 25-token RFE vs a 800-token doc
            # gives Jaccard ≈ 0.03 even with 100% coverage. Containment fixes this.
            intersection = len(tokens & doc_tokens)
            containment = intersection / max(len(tokens), 1)
            # Recency decay — prefer page's own doc_date over import date (critical for Confluence)
            eff_weight = base_weight
            if cfg.get('recency_decay', True):
                try:
                    date_str = (doc.get('doc_date') or doc.get('imported_at', ''))[:19]
                    if date_str:
                        page_date = datetime.fromisoformat(date_str)
                        if (datetime.utcnow() - page_date).days > 365:
                            eff_weight *= 0.8
                except: pass
            weighted = round(containment * eff_weight, 3)
            key = f"{doc_type}|{doc['title'][:40]}"
            breakdown[key] = max(breakdown.get(key, 0), weighted)
            if weighted > best_score:
                best_score = weighted; best_doc = doc

        if best_score >= min_score and best_doc:
            reason = f"Matched '{best_doc['title'][:50]}' (signal {best_score:.2f})"
            return status_label, reason, breakdown

    return 'no_match', 'No match found in any context source', breakdown

# ─── Confidence sentence ──────────────────────────────────────────────────────

def generate_confidence_sentence(cluster_title: str, match_status: str,
                                  match_reason: str, score: float) -> str:
    """Generate a one-sentence readable explanation of the cluster score."""
    doc_title = ""
    if match_reason:
        m = re.match(r"Matched '([^']+)'", match_reason)
        if m:
            doc_title = m.group(1)

    strength = "strongly" if score >= 0.75 else "partially" if score >= 0.45 else "weakly"

    if match_status == 'delivered':
        src = f" against {doc_title}" if doc_title else " against Release Notes"
        return (f"Matched {strength}{src} — this capability has already shipped "
                f"as part of a recent product release.")
    elif match_status == 'in_current_pi':
        if doc_title:
            return (f"Matched an active ADO epic ({doc_title}) — "
                    f"this feature is actively being built in the current PI.")
        return "Matched an active ADO epic — this feature is being built in the current PI."
    elif match_status == 'planned':
        if doc_title:
            return (f"Partial match found in {doc_title} — "
                    f"this capability appears to be planned for a future release.")
        return ("Matched a planned roadmap or backlog item — "
                "this capability is on the product roadmap.")
    return ""

# ─── Scoring pipeline ─────────────────────────────────────────────────────────

def _auto_migrate_release_notes(conn):
    """Before scoring: split any combined release-notes doc into per-version chunks.
    A combined doc has a title like 'Release Notes V4.32-V4.25.docx' — not a bare version string.
    Deduplicates chunks across multiple uploads of the same file (keeps the longest content
    per version). Runs automatically on every score so users don't need to re-upload.
    """
    cur = conn.cursor()
    cur.execute("SELECT id, title, content FROM context_docs WHERE doc_type='release_notes'")
    docs = [dict(r) for r in cur.fetchall()]

    combined_ids = []
    version_chunks: Dict[str, str] = {}  # version_label -> longest content seen
    for doc in docs:
        title = doc['title'] or ''
        if re.match(r'^[Vv]\d+\.\d+', title):
            continue  # already a version chunk — leave alone
        combined_ids.append(doc['id'])
        chunks = _chunk_release_notes(doc['content'] or '', title)
        if len(chunks) <= 1:
            continue  # no version headings detected
        for ver_label, chunk_text in chunks:
            if ver_label not in version_chunks or len(chunk_text) > len(version_chunks[ver_label]):
                version_chunks[ver_label] = chunk_text

    if version_chunks:
        for doc_id in combined_ids:
            cur.execute("DELETE FROM context_docs WHERE id=?", (doc_id,))
        imported_at = datetime.utcnow().isoformat()
        for ver_label, chunk_text in version_chunks.items():
            cur.execute(
                "INSERT INTO context_docs (doc_type,title,content,imported_at) VALUES (?,?,?,?)",
                ('release_notes', ver_label, chunk_text[:80000], imported_at)
            )
        conn.commit()

# ── Domain taxonomy — defined here so score_run + import_csv can use it ───────

CANONICAL_DOMAINS = [
    "Endpoint", "Cloud", "SIEM", "Identity", "ESPM",
    "Automation", "Platform", "MSP", "Email", "Mobile", "On-prem", "Reports",
]

_DOMAIN_KEYWORDS: list = [
    (["endpoint", "edr", "epp", "agent", "windows", "linux", "mac", "desktop",
       "workstation", "vulnerability", "misconfiguration", "cve", "inventory",
       "application control", "defender"], "Endpoint"),
    (["cloud", "aws", "azure", "gcp", "cspm", "alibaba", "s3", "bucket",
       "saas", "sandblast", "checkpoint harmony"], "Cloud"),
    (["siem", "clm", "log", "ingestion", "splunk", "qradar", "sentinel",
       "libraesva", "email security", "armis", "mikrotik", "checkpoint",
       "data source", "indexing", "field event", "cold storage", "firewall"], "SIEM"),
    (["identity", "ad", "active directory", "ldap", "mfa", "okta",
       "saml", "sso", "privileged", "user account"], "Identity"),
    (["espm", "spm", "posture", "exposure", "risk score", "attack surface"], "ESPM"),
    (["automation", "playbook", "soar", "response", "action", "remediation",
       "workflow", "trigger", "script"], "Automation"),
    (["platform", "api", "integration", "webhook", "sdk", "tenant",
       "role", "permission", "operator", "read only", "rbac"], "Platform"),
    (["msp", "multi-tenant", "managed service", "site", "cynet operator"], "MSP"),
    (["email", "mail", "phishing", "spam", "smtp"], "Email"),
    (["mobile", "ios", "android", "phone", "tablet"], "Mobile"),
    (["on-prem", "on_prem", "onprem", "local", "offline"], "On-prem"),
    (["report", "dashboard", "analytic", "chart", "export", "pdf",
       "csv", "graph", "trend"], "Reports"),
]

def _map_domain(raw_domain: str, cluster_title: str = "") -> str:
    """Map a raw Salesforce domain + free text to one of the 12 canonical domains."""
    haystack = f"{raw_domain or ''} {cluster_title or ''}".lower()
    for keywords, canonical in _DOMAIN_KEYWORDS:
        if any(kw in haystack for kw in keywords):
            return canonical
    for d in CANONICAL_DOMAINS:
        if d.lower() == (raw_domain or "").lower():
            return d
    return raw_domain or "Other"

async def score_run(run_id: str, job_id: str):  # noqa: C901
    """
    New pipeline (v1.2):
      Phase 1 — Score every RFE individually against context docs.
      Phase 2 — Cluster only matched RFEs (demand enrichment).
      Phase 3 — Write cluster records; solo matched RFEs become 1-member clusters.
    """
    conn = get_db()
    _auto_migrate_release_notes(conn)
    rows = [dict(r) for r in conn.execute(
        "SELECT * FROM rfe_pulls WHERE run_id=?", (run_id,)
    ).fetchall()]
    ctx_docs = [dict(r) for r in conn.execute(
        "SELECT id,doc_type,title,content,imported_at,doc_date FROM context_docs"
    ).fetchall()]
    conn.close()

    if not rows:
        jobs[job_id].update(status="done", message="No RFEs to score.", cluster_count=0)
        return

    scoring_cfg  = get_scoring_config()
    threshold    = scoring_cfg['similarity_threshold']
    high_thr     = scoring_cfg.get('threshold_high', 0.75)
    med_thr      = scoring_cfg.get('threshold_medium', 0.45)
    total        = len(rows)

    # ── PHASE 1: Score every RFE individually against context docs ─────────────
    jobs[job_id].update(status="scoring",
                        message=f"Phase 1 of 2 — Scoring {total} RFEs against context… 0 / {total}",
                        total_pairs=total, processed_pairs=0)

    rfe_scores: Dict[int, dict] = {}   # rfe_id → {match_status, match_reason, ctx_breakdown, signal_score}
    matched_rows: list = []            # only RFEs that scored above any threshold

    for i, rfe in enumerate(rows):
        rfe_text = f"{rfe.get('subject','') or ''} {(rfe.get('description') or '')[:500]}"
        m_status, m_reason, m_bd = determine_match_status(rfe_text, ctx_docs, scoring_cfg)
        signal = max(m_bd.values()) if m_bd else 0.0
        rfe_scores[rfe["id"]] = {
            "match_status":  m_status,
            "match_reason":  m_reason,
            "ctx_breakdown": m_bd,
            "signal_score":  signal,
        }
        if m_status != "no_match":
            matched_rows.append(rfe)
        if (i + 1) % 100 == 0 or i == total - 1:
            jobs[job_id].update(
                processed_pairs=i + 1,
                message=(f"Phase 1 of 2 — Scoring RFEs… {i+1} / {total} "
                         f"({len(matched_rows)} matched so far)"),
            )

    # Persist how many were filtered out
    conn2 = get_db()
    conn2.execute("UPDATE run_meta SET filtered_rfe_count=? WHERE run_id=?",
                  (total - len(matched_rows), run_id))
    conn2.commit(); conn2.close()

    jobs[job_id].update(
        message=(f"Phase 1 complete — {len(matched_rows)} / {total} RFEs matched. "
                 f"Starting clustering…"),
    )

    # ── PHASE 2: Cluster only matched RFEs ────────────────────────────────────
    api_key   = os.getenv("ANTHROPIC_API_KEY", "").strip()
    use_claude = bool(api_key)
    client    = None
    if use_claude:
        try:
            import anthropic as _anthropic
            client = _anthropic.AsyncAnthropic(api_key=api_key)
        except ImportError:
            use_claude = False

    mode_label = "Claude AI" if use_claude else "text similarity"

    def _cluster_group_key(rfe):
        domain = (rfe.get("domain") or "Unknown").strip()
        sub    = (rfe.get("sub_domain") or "").strip()
        return f"{domain}|{sub}" if sub else domain

    domain_groups: Dict[str, list] = defaultdict(list)
    for rfe in matched_rows:
        domain_groups[_cluster_group_key(rfe)].append(rfe)

    rfe_index  = {rfe["id"]: i for i, rfe in enumerate(matched_rows)}
    uf         = UnionFind(len(matched_rows))
    pair_best:   Dict[int, float] = {}
    pair_reason: Dict[int, dict]  = {}

    total_pairs = sum(len(g) * (len(g) - 1) // 2
                      for g in domain_groups.values() if len(g) >= 2)
    processed = 0

    jobs[job_id].update(
        status="clustering",
        message=(f"Phase 2 of 2 — Clustering {len(matched_rows)} matched RFEs "
                 f"({total_pairs} pairs, {mode_label})… 0 / {total_pairs}"),
        total_pairs=total_pairs, processed_pairs=0, scoring_mode=mode_label,
    )

    import time as _time
    t_start = _time.time()

    for domain, group in domain_groups.items():
        if len(group) < 2:
            continue
        pairs = list(itertools.combinations(group, 2))
        if use_claude:
            batch_size = 5
            for i in range(0, len(pairs), batch_size):
                batch = pairs[i:i + batch_size]
                results = await asyncio.gather(
                    *[claude_similarity(client, r1, r2) for r1, r2 in batch],
                    return_exceptions=True)
                for (r1, r2), result in zip(batch, results):
                    if not isinstance(result, Exception):
                        _apply_pair(result, r1, r2, rfe_index, uf,
                                    pair_best, pair_reason, threshold)
                    processed += 1
                elapsed = _time.time() - t_start
                rate = processed / elapsed if elapsed > 0 else 1
                rem = int((total_pairs - processed) / rate) if rate > 0 else 0
                eta = f"~{rem//60}m {rem%60}s remaining" if rem > 10 else "almost done"
                jobs[job_id].update(processed_pairs=processed,
                                    message=f"Phase 2 of 2 — Clustering… {processed} / {total_pairs} — {eta}")
        else:
            for r1, r2 in pairs:
                result = basic_similarity(r1, r2)
                _apply_pair(result, r1, r2, rfe_index, uf,
                            pair_best, pair_reason, threshold)
                processed += 1
            elapsed = _time.time() - t_start
            rate = processed / elapsed if elapsed > 0 else 1
            rem = int((total_pairs - processed) / rate) if rate > 0 and processed < total_pairs else 0
            eta = f"~{rem//60}m {rem%60}s remaining" if rem > 10 else "almost done"
            jobs[job_id].update(processed_pairs=processed,
                                message=f"Phase 2 of 2 — Clustering… {processed} / {total_pairs} — {eta}")

    # Build Union-Find groups
    clusters_map: Dict[int, list] = defaultdict(list)
    for i, rfe in enumerate(matched_rows):
        clusters_map[uf.find(i)].append(rfe)

    # ── PHASE 3: Write cluster records ────────────────────────────────────────
    _STATUS_PRIORITY = ["delivered", "in_current_pi", "planned", "no_match"]
    def _status_rank(s): return _STATUS_PRIORITY.index(s) if s in _STATUS_PRIORITY else 99

    conn = get_db(); cur = conn.cursor()
    created_at    = datetime.utcnow().isoformat()
    cluster_count = 0

    for root, members in clusters_map.items():
        is_solo       = len(members) == 1
        sim_data      = pair_reason.get(root, {})
        raw_sim_score = pair_best.get(root, 0.0)

        # Confidence score & band
        if is_solo:
            conf_score = rfe_scores[members[0]["id"]]["signal_score"]
        else:
            conf_score = raw_sim_score
        band = ("high"   if conf_score >= high_thr else
                "medium" if conf_score >= med_thr  else "low")

        # Cluster title = subject of highest-ARR member
        best_member = max(members, key=lambda r: r.get("account_arr") or 0)
        title        = best_member.get("subject") or best_member.get("account_name") or "Untitled"
        member_ids   = [m["id"] for m in members]
        total_arr    = sum(m.get("account_arr") or 0 for m in members)
        unique_accts = len({m["account_name"] for m in members if m.get("account_name")})

        # Per-member match data (already computed in Phase 1)
        member_matches = [
            {
                "id":           m["id"],
                "match_status": rfe_scores[m["id"]]["match_status"],
                "ctx_breakdown":rfe_scores[m["id"]]["ctx_breakdown"],
            }
            for m in members
        ]

        # Cluster-level status = best individual status (delivered > in_pi > planned)
        best_mm    = min(member_matches, key=lambda x: _status_rank(x["match_status"]))
        match_status = best_mm["match_status"]

        # Cluster context breakdown = union of member breakdowns (highest score per key)
        ctx_breakdown: dict = {}
        for mm in member_matches:
            for k, v in mm["ctx_breakdown"].items():
                if v > ctx_breakdown.get(k, 0):
                    ctx_breakdown[k] = v

        match_reason = ""
        if ctx_breakdown:
            best_k   = max(ctx_breakdown, key=ctx_breakdown.get)
            doc_name = best_k.split("|")[-1]
            match_reason = f"Matched '{doc_name}' (signal {ctx_breakdown[best_k]:.2f})"

        score_breakdown = json.dumps({
            "similarity_reason": (sim_data.get("reason", "Individual context match")
                                  if is_solo else
                                  sim_data.get("reason", "Shared key terms in RFE descriptions")),
            "common_terms":      sim_data.get("common", []),
            "confidence":        round(conf_score, 3),
            "context_breakdown": ctx_breakdown,
            "match_score":       max(ctx_breakdown.values()) if ctx_breakdown else 0,
            "rfe_count":         len(members),
            "domains":           list({m.get("domain", "") for m in members if m.get("domain")}),
            "member_matches":    member_matches,
            "solo":              is_solo,
        })

        confidence_sentence = generate_confidence_sentence(
            title, match_status, match_reason, conf_score)

        cur.execute(
            "INSERT INTO rfe_clusters "
            "(run_id,cluster_title,confidence_score,confidence_band,"
            "member_case_numbers,total_arr,account_count,status,match_status,match_reason,"
            "score_breakdown,confidence_sentence,created_at) "
            "VALUES (?,?,?,?,?,?,?,'pending',?,?,?,?,?)",
            (run_id, title, conf_score, band, json.dumps(member_ids),
             total_arr, unique_accts,
             match_status, match_reason, score_breakdown,
             confidence_sentence, created_at),
        )
        cluster_count += 1

    cur.execute(
        "UPDATE run_meta SET cluster_count=?,status='scored',completed_at=? WHERE run_id=?",
        (cluster_count, created_at, run_id),
    )
    conn.commit(); conn.close()
    jobs[job_id].update(
        status="done",
        message=(f"Done — {cluster_count} signals found "
                 f"({len(matched_rows)} matched RFEs, {mode_label})."),
        cluster_count=cluster_count,
    )

def _apply_pair(result, r1, r2, rfe_index, uf, pair_best, pair_reason, threshold):
    score = float(result.get("score", 0))

    # Fix 2 — domain/sub-domain bonus: same taxonomy = easier to cluster.
    # Applied only for the threshold gate; raw score is stored for confidence band.
    d1 = (r1.get("domain") or "").strip().lower()
    d2 = (r2.get("domain") or "").strip().lower()
    s1 = (r1.get("sub_domain") or "").strip().lower()
    s2 = (r2.get("sub_domain") or "").strip().lower()
    domain_bonus = 0.0
    if d1 and d1 == d2:
        domain_bonus += 0.10          # same domain
        if s1 and s1 == s2:
            domain_bonus += 0.10      # same domain + same sub-domain → 0.20 total

    if score + domain_bonus >= threshold:
        idx1, idx2 = rfe_index[r1["id"]], rfe_index[r2["id"]]
        uf.union(idx1, idx2)
        root = uf.find(idx1)
        if pair_best.get(root, 0) < score:   # store raw score (not inflated)
            pair_best[root] = score
            pair_reason[root] = result

# ─── CSV Import ───────────────────────────────────────────────────────────────

@app.post("/api/import-csv")
async def import_csv(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    content = await file.read()
    try: text = content.decode("utf-8-sig")
    except UnicodeDecodeError: text = content.decode("latin-1")

    reader = csv.DictReader(io.StringIO(text))
    headers = reader.fieldnames or []
    col_map = detect_columns(list(headers))
    run_id = str(uuid.uuid4()); job_id = str(uuid.uuid4())
    pulled_at = datetime.utcnow().isoformat()
    jobs[job_id] = {"status": "importing", "message": "Parsing CSV…", "run_id": run_id, "job_id": job_id}
    records = list(reader)
    conn = get_db(); cur = conn.cursor()
    for rec in records:
        def g(field):
            col = col_map.get(field)
            return rec.get(col, "").strip() if col else ""
        cur.execute(
            "INSERT INTO rfe_pulls (run_id,case_number,subject,description,account_name,account_arr,"
            "status,domain,sub_domain,severity,created_date,pulled_at) VALUES (?,?,?,?,?,?,?,?,?,?,?,?)",
            (run_id, g("case_number"), g("subject"), g("description"), g("account_name"),
             parse_arr(g("account_arr")), g("status"), g("domain"), g("sub_domain"),
             g("severity"), g("created_date"), pulled_at),
        )
    cur.execute("INSERT OR REPLACE INTO run_meta (run_id,started_at,rfe_count,status,source) VALUES (?,?,?,'imported','csv')",
                (run_id, pulled_at, len(records)))
    conn.commit()

    # Fix 3 — Infer domain for blank-domain RFEs before scoring begins.
    # RFEs without a domain field land in a useless mixed bucket; keyword inference
    # puts them next to the right peers (CLM→SIEM, S3→Cloud, etc.).
    blank_rows = conn.execute(
        "SELECT id, subject, description FROM rfe_pulls "
        "WHERE run_id=? AND (domain IS NULL OR domain='')", (run_id,)
    ).fetchall()
    inferred_count = 0
    for br in blank_rows:
        text = f"{br['subject'] or ''} {br['description'] or ''}"
        inferred = _map_domain("", text)
        if inferred and inferred != "Other":
            conn.execute("UPDATE rfe_pulls SET domain=? WHERE id=?", (inferred, br["id"]))
            inferred_count += 1
    if inferred_count:
        conn.commit()

    conn.close()
    jobs[job_id].update(status="imported",
                        message=f"Imported {len(records)} RFEs ({inferred_count} domains inferred) — starting scoring…",
                        rfe_count=len(records))
    background_tasks.add_task(score_run, run_id, job_id)
    return {"job_id": job_id, "run_id": run_id, "rfe_count": len(records), "columns_detected": col_map}

# ─── Context docs ─────────────────────────────────────────────────────────────

# Version-heading pattern. Matches the various heading shapes that appear in
# real release notes files, e.g.
#   "Release Notes V4.32 (SaaS)"
#   "xRelease Notes V4.32 (SaaS)"  (leading x is a docx-extraction artifact)
#   "## V4.32"
#   "Version 4.32"
#   "V4.32"
# Rejects lines that have extra text after the version (paragraphs, sentences).
_VERSION_HEADING = re.compile(
    r'^\s*(?:#{1,4}\s*)?'                              # optional markdown header markers
    r'x?'                                              # optional leading x (docx artifact)
    r'(?:Release\s+Notes\s+|Version\s+|Release\s+)?'   # optional prefix words
    r'[Vv]?(\d{1,2}\.\d{1,2}(?:\.\d{1,2})?)'           # version digits 4.32 or 4.32.1
    r'(?:\s*\([^)\n]{1,40}\))?'                        # optional parenthetical e.g. (SaaS)
    r'\s*$',                                           # end of line — no trailing text
    re.MULTILINE | re.IGNORECASE
)

def _chunk_release_notes(text: str, filename: str):
    """Split release notes into per-version chunks.
    Returns list of (version_label, chunk_text) e.g. [('V4.32', '...'), ('V4.31', '...')].
    Falls back to a single chunk with the filename as title if no version headings found.
    """
    matches = list(_VERSION_HEADING.finditer(text))
    if not matches:
        return [(filename, text)]
    chunks = []
    for i, m in enumerate(matches):
        ver_label = f"V{m.group(1)}"
        start = m.start()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append((ver_label, chunk))
    return chunks if chunks else [(filename, text)]

@app.post("/api/context/release-notes")
async def import_release_notes(file: UploadFile = File(...)):
    content = await file.read()
    filename = file.filename or ""

    if filename.lower().endswith(".docx"):
        try:
            import docx as _docx, io as _io
            doc = _docx.Document(_io.BytesIO(content))
            # Chunk by Word Heading 1 — each heading marks a version boundary
            chunks = []
            current_heading = None
            current_parts = []
            for p in doc.paragraphs:
                style_name = (p.style.name or '').lower()
                style_id = (p.style.style_id or '').lower().replace(' ', '')
                is_h1 = 'heading 1' in style_name or style_id == 'heading1'
                if is_h1 and p.text.strip():
                    if current_heading is not None:
                        section_text = '\n'.join(current_parts)
                        if section_text.strip():
                            chunks.append((current_heading, section_text))
                    ver_match = re.search(r'[Vv]?(\d{1,2}\.\d{1,2}(?:\.\d{1,2})?)', p.text)
                    current_heading = f"V{ver_match.group(1)}" if ver_match else p.text.strip()[:40]
                    current_parts = [p.text.strip()]
                elif p.text.strip():
                    current_parts.append(p.text.strip())
            if current_heading is not None:
                section_text = '\n'.join(current_parts)
                if section_text.strip():
                    chunks.append((current_heading, section_text))
            # For chars reporting + fallback
            text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
            # Fallback: if no Heading 1 found, use text-based regex chunker
            if not chunks:
                chunks = _chunk_release_notes(text, filename)
        except Exception as e:
            raise HTTPException(400, f"Could not read .docx: {e}")
    else:
        try: text = content.decode("utf-8-sig")
        except UnicodeDecodeError: text = content.decode("latin-1")
        chunks = _chunk_release_notes(text, filename)

    conn = get_db()
    imported_at = datetime.utcnow().isoformat()
    # Clean replace: drop existing release_notes docs so re-uploading is idempotent
    conn.execute("DELETE FROM context_docs WHERE doc_type='release_notes'")
    for ver_label, chunk_text in chunks:
        conn.execute("INSERT INTO context_docs (doc_type,title,content,imported_at) VALUES ('release_notes',?,?,?)",
                     (ver_label, chunk_text[:80000], imported_at))
    conn.commit(); conn.close()
    return {"imported": True, "title": filename, "chars": len(text), "chunks": len(chunks)}

@app.post("/api/context/ado-csv")
async def import_ado_csv(file: UploadFile = File(...), sub_type: str = Form("ado_backlog")):
    """Import ADO Epics CSV. sub_type: 'ado_current_pi' or 'ado_backlog'
    Each epic row is stored as an individual context_doc so per-RFE scoring can reference specific Epic IDs.
    """
    if sub_type not in ("ado_current_pi", "ado_backlog"):
        sub_type = "ado_backlog"
    content = await file.read()
    try: text = content.decode("utf-8-sig")
    except UnicodeDecodeError: text = content.decode("latin-1")
    reader = csv.DictReader(io.StringIO(text))
    records = list(reader)
    conn = get_db()
    imported_at = datetime.utcnow().isoformat()
    # Clean replace: drop existing docs of this sub_type so re-uploading is idempotent
    conn.execute("DELETE FROM context_docs WHERE doc_type=?", (sub_type,))
    inserted = 0
    for i, rec in enumerate(records[:500]):
        # Detect Epic ID from common ADO column names
        epic_id_raw = (rec.get('ID') or rec.get('Work Item Id') or rec.get('Work Item ID') or
                       rec.get('id') or rec.get('WorkItemId') or '').strip()
        epic_title = (rec.get('Title') or rec.get('title') or rec.get('Name') or '').strip()
        # Build searchable content from all relevant fields
        parts = []
        for k, v in rec.items():
            if v and any(kw in k.lower() for kw in
                         ['title','name','description','summary','acceptance','tags','area','iteration']):
                clean = _strip_html(str(v).strip())
                if clean:
                    parts.append(clean)
        content_str = " | ".join(parts)
        if not content_str:
            continue
        # Display title embeds the Epic ID so it can be copied from the UI
        if epic_id_raw and epic_title:
            display_title = f"{epic_title} #{epic_id_raw}"
        elif epic_id_raw:
            display_title = f"Epic #{epic_id_raw}"
        elif epic_title:
            display_title = epic_title
        else:
            display_title = f"Epic {i+1}"
        conn.execute("INSERT INTO context_docs (doc_type,title,content,imported_at) VALUES (?,?,?,?)",
                     (sub_type, display_title, content_str[:5000], imported_at))
        inserted += 1
    conn.commit(); conn.close()
    return {"imported": True, "rows": len(records), "inserted": inserted, "doc_type": sub_type}

@app.post("/api/context/backlog-csv")
async def import_backlog(file: UploadFile = File(...)):
    content = await file.read()
    try: text = content.decode("utf-8-sig")
    except UnicodeDecodeError: text = content.decode("latin-1")
    reader = csv.DictReader(io.StringIO(text))
    records = list(reader)
    summary = json.dumps(records[:500], ensure_ascii=False)
    conn = get_db()
    conn.execute("INSERT INTO context_docs (doc_type,title,content,imported_at) VALUES ('ado_backlog',?,?,?)",
                 (file.filename, summary, datetime.utcnow().isoformat()))
    conn.commit(); conn.close()
    return {"imported": True, "title": file.filename, "rows": len(records)}

@app.post("/api/context/roadmap")
async def import_roadmap(file: UploadFile = File(...)):
    """Import PPTX or PDF roadmap file."""
    content = await file.read()
    filename = file.filename or ""
    text = ""

    if filename.lower().endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            text = "\n".join(parts)
        except ImportError:
            raise HTTPException(400, "python-pptx not installed. Run: pip install python-pptx")
        except Exception as e:
            raise HTTPException(400, f"Could not read PPTX: {e}")
    elif filename.lower().endswith(".pdf"):
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(content))
            parts = [page.extract_text() or "" for page in reader.pages]
            text = "\n".join(parts)
        except ImportError:
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(content))
                parts = [page.extract_text() or "" for page in reader.pages]
                text = "\n".join(parts)
            except ImportError:
                raise HTTPException(400, "pypdf not installed. Run: pip install pypdf")
        except Exception as e:
            raise HTTPException(400, f"Could not read PDF: {e}")
    else:
        try: text = content.decode("utf-8-sig")
        except UnicodeDecodeError: text = content.decode("latin-1")

    conn = get_db()
    conn.execute("INSERT INTO context_docs (doc_type,title,content,imported_at) VALUES ('roadmap',?,?,?)",
                 (filename, text[:50000], datetime.utcnow().isoformat()))
    conn.commit(); conn.close()
    return {"imported": True, "title": filename, "chars": len(text)}

# ─── Confluence ───────────────────────────────────────────────────────────────

def _cf_creds():
    base = os.getenv("CONFLUENCE_BASE_URL", "").rstrip("/")
    email = os.getenv("CONFLUENCE_EMAIL", "") or os.getenv("CONFLUENCE_USERNAME", "")
    token = os.getenv("CONFLUENCE_API_TOKEN", "")
    return base, email, token

def _cf_get(path: str) -> dict:
    import urllib.request, urllib.error, base64
    base, email, token = _cf_creds()
    url = f"{base}/wiki/rest/api{path}"
    creds = base64.b64encode(f"{email}:{token}".encode()).decode()
    req = urllib.request.Request(url, headers={"Authorization": f"Basic {creds}", "Accept": "application/json"})
    try:
        with urllib.request.urlopen(req, timeout=15) as resp:
            return json.loads(resp.read())
    except urllib.error.HTTPError as e:
        if e.code == 401: raise Exception("401 Unauthorized — check CONFLUENCE_EMAIL and CONFLUENCE_API_TOKEN")
        if e.code == 403: raise Exception("403 Forbidden")
        raise Exception(f"HTTP {e.code}: {e.reason}")

def _strip_html(html: str) -> str:
    text = re.sub(r"<[^>]+>", " ", html or "")
    return re.sub(r"\s+", " ", text).strip()

@app.get("/api/confluence/status")
async def confluence_status():
    base, email, token = _cf_creds()
    return {"configured": bool(base and email and token), "base_url": base, "email": email}

@app.get("/api/confluence/spaces")
async def confluence_spaces():
    base, email, token = _cf_creds()
    if not base or not token:
        raise HTTPException(400, "Confluence credentials not set in .env")
    loop = asyncio.get_event_loop()
    try:
        data = await loop.run_in_executor(executor, _cf_get, "/space?limit=50&type=global")
        spaces = [{"key": s["key"], "name": s["name"],
                   "description": s.get("description", {}).get("plain", {}).get("value", "")}
                  for s in data.get("results", [])]
        return {"spaces": spaces}
    except Exception as e:
        raise HTTPException(400, f"Could not connect to Confluence: {e}. "
                                 f"Check CONFLUENCE_BASE_URL={base} and CONFLUENCE_EMAIL={email}")

@app.get("/api/confluence/page-tree")
async def confluence_page_tree(space_key: str):
    """Return flat page list with parent info so UI can build a tree."""
    loop = asyncio.get_event_loop()
    pages = []
    start = 0; limit = 100
    try:
        while True:
            data = await loop.run_in_executor(executor, _cf_get,
                f"/content?spaceKey={space_key}&type=page&status=current"
                f"&expand=ancestors,version&limit={limit}&start={start}")
            batch = data.get("results", [])
            if not batch: break
            base, _, _ = _cf_creds()
            for p in batch:
                ancestors = p.get("ancestors", [])
                parent_id = ancestors[-1]["id"] if ancestors else None
                last_modified = p.get("version", {}).get("when", "")
                pages.append({
                    "id": p["id"], "title": p["title"], "parent_id": parent_id,
                    "last_modified": last_modified,
                    "url": f"{base}/wiki/spaces/{space_key}/pages/{p['id']}"
                })
            if len(batch) < limit: break
            start += limit
        return {"pages": pages, "total": len(pages)}
    except Exception as e:
        raise HTTPException(400, str(e))

class ImportPagesRequest(BaseModel):
    space_key: str
    page_ids: List[str]

@app.post("/api/confluence/import-pages")
async def import_confluence_pages(background_tasks: BackgroundTasks, req: ImportPagesRequest):
    job_id = str(uuid.uuid4())
    jobs[job_id] = {"status": "running", "message": "Starting page import…",
                    "imported": 0, "total": len(req.page_ids)}
    background_tasks.add_task(_import_pages_task, req.space_key, req.page_ids, job_id)
    return {"job_id": job_id}

async def _import_pages_task(space_key: str, page_ids: List[str], job_id: str):
    loop = asyncio.get_event_loop()
    conn = get_db()
    imported = 0
    base, _, _ = _cf_creds()
    try:
        for page_id in page_ids:
            try:
                data = await loop.run_in_executor(executor, _cf_get,
                    f"/content/{page_id}?expand=body.storage,version")
                title = data.get("title", "Untitled")
                body_html = data.get("body", {}).get("storage", {}).get("value", "")
                clean = _strip_html(body_html)
                url = f"{base}/wiki/spaces/{space_key}/pages/{page_id}"
                # Use page's own last-modified date for recency scoring
                doc_date = data.get("version", {}).get("when", "")
                conn.execute(
                    "INSERT INTO context_docs (doc_type,title,content,source_url,imported_at,doc_date) VALUES ('confluence',?,?,?,?,?)",
                    (f"[{space_key}] {title}", clean[:30000], url, datetime.utcnow().isoformat(), doc_date)
                )
                imported += 1
                jobs[job_id].update(imported=imported,
                                    message=f"Imported {imported}/{len(page_ids)} pages…")
            except Exception:
                pass  # skip failed pages
        conn.commit()
        jobs[job_id].update(status="done", message=f"✅ Imported {imported} pages")
    except Exception as e:
        jobs[job_id].update(status="error", message=str(e))
    finally:
        conn.close()

@app.post("/api/confluence/import-space")
async def import_confluence_space(background_tasks: BackgroundTasks,
                                   space_key: str = Form(...), space_name: str = Form("")):
    job_id = str(uuid.uuid4())
    jobs[job_id] = {"status": "running", "message": "Starting Confluence import…",
                    "imported": 0, "space": space_name or space_key}
    background_tasks.add_task(_import_space_task, space_key, space_name, job_id)
    return {"job_id": job_id}

async def _import_space_task(space_key: str, space_name: str, job_id: str):
    loop = asyncio.get_event_loop()
    try:
        start = 0; limit = 25; imported = 0
        conn = get_db()
        while True:
            data = await loop.run_in_executor(executor, _cf_get,
                f"/content?spaceKey={space_key}&type=page&status=current"
                f"&expand=body.storage,version&limit={limit}&start={start}")
            pages = data.get("results", [])
            if not pages: break
            for page in pages:
                title = page.get("title", "Untitled")
                body_html = page.get("body", {}).get("storage", {}).get("value", "")
                clean = _strip_html(body_html)
                page_id = page.get("id", "")
                url = f"{os.getenv('CONFLUENCE_BASE_URL','').rstrip('/')}/wiki/spaces/{space_key}/pages/{page_id}"
                doc_date = page.get("version", {}).get("when", "")
                conn.execute(
                    "INSERT INTO context_docs (doc_type,title,content,source_url,imported_at,doc_date) VALUES ('confluence',?,?,?,?,?)",
                    (f"[{space_name or space_key}] {title}", clean[:30000], url, datetime.utcnow().isoformat(), doc_date),
                )
                imported += 1
                jobs[job_id].update(imported=imported, message=f"Imported {imported} pages from {space_name or space_key}…")
            conn.commit()
            if len(pages) < limit: break
            start += limit
        conn.close()
        jobs[job_id].update(status="done", message=f"✅ Imported {imported} pages from {space_name or space_key}")
    except Exception as e:
        jobs[job_id].update(status="error", message=str(e))

@app.get("/api/context")
async def get_context_docs():
    conn = get_db()
    docs = [dict(r) for r in conn.execute(
        "SELECT id,doc_type,title,source_url,imported_at,length(content) as chars FROM context_docs ORDER BY imported_at DESC"
    ).fetchall()]
    conn.close()
    return {"docs": docs}

@app.delete("/api/context/{doc_id}")
async def delete_context_doc(doc_id: int):
    conn = get_db()
    conn.execute("DELETE FROM context_docs WHERE id=?", (doc_id,))
    conn.commit(); conn.close()
    return {"deleted": True}

# ─── Scoring Config endpoints ─────────────────────────────────────────────────

@app.get("/api/scoring-config")
async def get_scoring_config_api():
    return get_scoring_config()

class ScoringConfigUpdate(BaseModel):
    weight_release_notes: Optional[float] = None
    weight_ado_current_pi: Optional[float] = None
    weight_ado_backlog: Optional[float] = None
    weight_confluence_dated: Optional[float] = None
    weight_roadmap: Optional[float] = None
    weight_confluence_undated: Optional[float] = None
    recency_decay: Optional[bool] = None
    threshold_high: Optional[float] = None
    threshold_medium: Optional[float] = None
    similarity_threshold: Optional[float] = None
    high_confidence_threshold: Optional[float] = None
    match_thresh_delivered: Optional[float] = None
    match_thresh_in_pi: Optional[float] = None
    match_thresh_planned: Optional[float] = None

@app.put("/api/scoring-config")
async def update_scoring_config(cfg: ScoringConfigUpdate):
    data = cfg.model_dump(exclude_none=True)
    for key, val in data.items():
        set_cfg(key, str(val).lower() if isinstance(val, bool) else str(val))
    return {"saved": True}

# ─── Re-score existing run ────────────────────────────────────────────────────

@app.post("/api/rescore")
async def api_rescore(background_tasks: BackgroundTasks):
    """Re-run scoring on the most recent run's RFEs against current context docs."""
    conn = get_db(); cur = conn.cursor()
    cur.execute("SELECT run_id FROM run_meta ORDER BY started_at DESC LIMIT 1")
    row = cur.fetchone()
    if not row:
        conn.close(); raise HTTPException(400, "No existing run to re-score")
    run_id = row["run_id"]
    # Delete old clusters for this run so score_run creates fresh ones
    cur.execute("DELETE FROM rfe_clusters WHERE run_id=?", (run_id,))
    cur.execute("UPDATE run_meta SET status='imported' WHERE run_id=?", (run_id,))
    conn.commit(); conn.close()
    job_id = str(uuid.uuid4())
    jobs[job_id] = {"status": "scoring", "message": "Re-scoring against updated context sources…",
                    "run_id": run_id, "job_id": job_id}
    background_tasks.add_task(score_run, run_id, job_id)
    return {"job_id": job_id, "run_id": run_id}

# ─── Clusters ─────────────────────────────────────────────────────────────────

@app.get("/api/status/{job_id}")
async def api_status(job_id: str):
    job = jobs.get(job_id)
    if not job: raise HTTPException(404, "Job not found")
    return job

@app.get("/api/clusters")
async def api_clusters(run_id: Optional[str] = None):
    conn = get_db(); cur = conn.cursor()
    if not run_id:
        cur.execute("SELECT run_id FROM run_meta ORDER BY started_at DESC LIMIT 1")
        row = cur.fetchone()
        if not row: conn.close(); return {"clusters": [], "run_id": None}
        run_id = row["run_id"]
    cur.execute("SELECT * FROM rfe_clusters WHERE run_id=? ORDER BY confidence_score DESC", (run_id,))
    clusters = [dict(r) for r in cur.fetchall()]
    for c in clusters:
        member_ids = json.loads(c.get("member_case_numbers") or "[]")
        if member_ids:
            ph = ",".join("?" * len(member_ids))
            cur.execute(f"SELECT id,case_number,subject,description,account_name,account_arr,created_date,status,domain "
                        f"FROM rfe_pulls WHERE id IN ({ph})", member_ids)
            c["members"] = [dict(r) for r in cur.fetchall()]
        else:
            c["members"] = []
        # Parse score_breakdown JSON for UI
        try: c["score_breakdown"] = json.loads(c.get("score_breakdown") or "{}")
        except: c["score_breakdown"] = {}
        # Attach per-RFE match data to each member object
        member_matches_by_id = {mm["id"]: mm for mm in c["score_breakdown"].get("member_matches", [])}
        for member in c["members"]:
            mm = member_matches_by_id.get(member["id"])
            if mm:
                member["match_status"] = mm.get("match_status", "no_match")
                member["ctx_breakdown"] = mm.get("ctx_breakdown", {})
        # Ensure confidence_sentence is present
        if not c.get("confidence_sentence"):
            c["confidence_sentence"] = generate_confidence_sentence(
                c.get("cluster_title",""), c.get("match_status","no_match"),
                c.get("match_reason",""), c.get("confidence_score",0)
            )
    conn.close()
    return {"clusters": clusters, "run_id": run_id}

@app.get("/api/search-clusters")
async def api_search_clusters(q: str = ""):
    """Search clusters by Case ID, subject, keyword — returns rendered HTML fragment."""
    conn = get_db()
    run_row = conn.execute(
        "SELECT run_id FROM run_meta ORDER BY started_at DESC LIMIT 1"
    ).fetchone()
    if not run_row:
        conn.close()
        return HTMLResponse('<div class="empty-state"><div class="big-icon">📡</div>'
                            '<h3>No data yet</h3><p>Upload a CSV first.</p></div>')
    html, _ = _build_clusters_for_page(conn, run_row["run_id"], search_query=q)
    conn.close()
    return HTMLResponse(html)

@app.patch("/api/clusters/{cluster_id}/approve")
async def api_approve(cluster_id: int):
    conn = get_db(); cur = conn.cursor()
    cur.execute("UPDATE rfe_clusters SET status='approved' WHERE id=?", (cluster_id,))
    if cur.rowcount == 0: conn.close(); raise HTTPException(404, "Not found")
    conn.commit(); conn.close(); return {"status": "approved"}

@app.patch("/api/clusters/{cluster_id}/reject")
async def api_reject(cluster_id: int):
    conn = get_db()
    conn.execute("UPDATE rfe_clusters SET status='rejected' WHERE id=?", (cluster_id,))
    conn.commit(); conn.close(); return {"status": "rejected"}

@app.patch("/api/clusters/{cluster_id}/restore")
async def api_restore(cluster_id: int):
    conn = get_db()
    conn.execute("UPDATE rfe_clusters SET status='pending' WHERE id=?", (cluster_id,))
    conn.commit(); conn.close(); return {"status": "pending"}

@app.post("/api/clusters/approve-bulk-high")
async def api_approve_bulk_high(run_id: Optional[str] = None):
    conn = get_db(); cur = conn.cursor()
    if not run_id:
        cur.execute("SELECT run_id FROM run_meta ORDER BY started_at DESC LIMIT 1")
        row = cur.fetchone(); run_id = row["run_id"] if row else None
    if not run_id: conn.close(); return {"approved": 0}
    ht = float(get_cfg("HIGH_CONFIDENCE_THRESHOLD", "0.85"))
    cur.execute("UPDATE rfe_clusters SET status='approved' WHERE run_id=? AND confidence_score>=? AND status='pending'",
                (run_id, ht))
    approved = cur.rowcount; conn.commit(); conn.close()
    return {"approved": approved}

# ─── Email ────────────────────────────────────────────────────────────────────

class PreviewEmailRequest(BaseModel):
    cluster_id: int
    account_name: Optional[str] = None
    member_id: Optional[int] = None     # specific RFE ID selected in the drawer
    force_status: Optional[str] = None  # legacy / fallback override

async def _generate_email(account_name, case_number, subject, similar_count,
                           match_status="no_match", member_details=None) -> dict:
    """
    member_details: optional list of {case_number, subject, match_status} for accounts with
    multiple RFEs in the same cluster (potentially with different statuses).
    """
    api_key = os.getenv("ANTHROPIC_API_KEY", "").strip()

    _status_labels = {
        "delivered":     "already shipped in a recent release",
        "in_current_pi": "actively being built in the current sprint/PI",
        "planned":       "on the roadmap and planned for a future release",
        "no_match":      "being reviewed by the product team",
    }
    _status_icons = {
        "delivered":     "✅ Delivered in a recent release",
        "in_current_pi": "🔨 Actively being built in the current PI",
        "planned":       "📋 Planned for an upcoming release",
        "no_match":      "📝 Under review by our product team",
    }

    # Decide whether we have a mixed-status scenario
    is_multi = bool(member_details and len(member_details) > 1)
    has_mixed = is_multi and len({m["match_status"] for m in (member_details or [])}) > 1

    if api_key:
        import anthropic as _anthropic
        client = _anthropic.AsyncAnthropic(api_key=api_key)
        if has_mixed:
            rfe_lines = "\n".join(
                f"- Case {m['case_number']}: {m['subject']} → {_status_labels.get(m['match_status'], 'under review')}"
                for m in member_details
            )
            status_ctx = f"multiple related feature requests with different statuses:\n{rfe_lines}"
        else:
            status_ctx = _status_labels.get(match_status, "being reviewed")
        prompt = (
            "You are a product manager at a B2B security company called Cynet.\n"
            "Write a professional, warm, concise email to a customer about their feature request(s).\n"
            f"Status: {status_ctx}.\n"
            "Do NOT promise specific delivery dates. Keep it under 150 words.\n"
            f"Customer: {account_name} | Primary case: {case_number} | Topic: {subject}\n"
            f"Similar requests from other customers: {similar_count}\n\n"
            "Return JSON with: to, subject, body"
        )
        resp = await client.messages.create(model="claude-opus-4-5", max_tokens=600,
                                             messages=[{"role": "user", "content": prompt}])
        text = resp.content[0].text.strip()
        if "```" in text:
            text = text.split("```")[1]
            if text.startswith("json"): text = text[4:]
        return json.loads(text.strip())
    else:
        if has_mixed:
            # One email, one bullet per RFE with its own status
            bullets = "\n".join(
                f"  • Case {m['case_number']} – {m['subject']}\n"
                f"    {_status_icons.get(m['match_status'], '📝 Under review')}"
                for m in member_details
            )
            body = (
                f"Dear {account_name} team,\n\n"
                f"Thank you for your feature requests related to {subject}. "
                f"Here is the latest status on each:\n\n"
                f"{bullets}\n\n"
                f"You're not alone — {similar_count} similar requests from other customers have "
                f"been logged, which helps us prioritize effectively.\n\n"
                f"We appreciate your continued feedback and will keep you updated.\n\n"
                f"Best regards,\nCynet Product Team"
            )
            email_subject = f"Re: Your Feature Requests – {account_name}"
        else:
            status_line = {
                "delivered":     "great news — this feature has already been shipped in a recent release!",
                "in_current_pi": "this feature is actively being built in our current development cycle.",
                "planned":       "this feature is on our roadmap and planned for an upcoming release.",
                "no_match":      "your request is on our radar and actively being reviewed by the product team.",
            }.get(match_status, "your request is being reviewed.")
            body = (
                f"Dear {account_name} team,\n\nThank you for submitting your feature request "
                f"(Case {case_number}: {subject}).\n\nWe have {status_line}\n\n"
                f"You're not alone — {similar_count} similar requests from other customers have been logged, "
                f"which helps us prioritize effectively.\n\n"
                f"We appreciate your feedback and will keep you updated as things progress.\n\n"
                f"Best regards,\nCynet Product Team"
            )
            email_subject = f"Re: Your Feature Request – Case {case_number}"
        return {
            "to": f"contact@{account_name.lower().replace(' ','')}.com",
            "subject": email_subject,
            "body": body,
            "_template": True
        }

@app.post("/api/preview-email")
async def api_preview_email(req: PreviewEmailRequest):
    conn = get_db(); cur = conn.cursor()
    cur.execute("SELECT * FROM rfe_clusters WHERE id=?", (req.cluster_id,))
    cluster = dict(cur.fetchone() or {})
    if not cluster: conn.close(); raise HTTPException(404, "Cluster not found")
    member_ids = json.loads(cluster.get("member_case_numbers") or "[]")

    # Per-RFE status lookup table
    try:
        sb = json.loads(cluster.get("score_breakdown") or "{}")
        mm_by_id = {mm["id"]: mm for mm in sb.get("member_matches", [])}
    except Exception:
        mm_by_id = {}

    def _member_status(mid):
        mm = mm_by_id.get(mid)
        return (mm.get("match_status") if mm and mm.get("match_status") else None) \
               or cluster.get("match_status", "no_match")

    # ── Fast path: specific RFE selected in the drawer ────────────────────────
    if req.member_id:
        cur.execute("SELECT id,case_number,subject,account_name FROM rfe_pulls WHERE id=?",
                    (req.member_id,))
        row = cur.fetchone()
        conn.close()
        if not row:
            raise HTTPException(404, "RFE not found")
        m = dict(row)
        effective_status = _member_status(req.member_id)
        result = await _generate_email(
            m.get("account_name", "Customer"),
            m.get("case_number", "N/A"),
            m.get("subject") or cluster.get("cluster_title", "Feature Request"),
            len(member_ids),
            effective_status,
            member_details=None,
        )
        result["effective_status"] = effective_status
        result["rfe_subject"] = m.get("subject", "")
        return result

    # ── Fallback: no member_id — find all members for the account ─────────────
    account_members = []
    if member_ids:
        ph = ",".join("?" * len(member_ids))
        if req.account_name:
            cur.execute(
                f"SELECT id,case_number,subject,account_name FROM rfe_pulls "
                f"WHERE id IN ({ph}) AND account_name=?",
                member_ids + [req.account_name]
            )
            account_members = [dict(r) for r in cur.fetchall()]
        if not account_members:
            cur.execute(f"SELECT id,case_number,subject,account_name FROM rfe_pulls WHERE id=? LIMIT 1",
                        (member_ids[0],))
            row = cur.fetchone()
            if row: account_members = [dict(row)]
    conn.close()

    primary = account_members[0] if account_members else {}
    member_details = [{
        "case_number":  m.get("case_number", "N/A"),
        "subject":      m.get("subject", ""),
        "match_status": _member_status(m["id"]),
    } for m in account_members]

    _priority = ["delivered", "in_current_pi", "planned", "no_match"]
    effective_status = min(
        (d["match_status"] for d in member_details),
        key=lambda s: _priority.index(s) if s in _priority else 99,
        default=cluster.get("match_status", "no_match"),
    )
    if req.force_status and req.force_status in ('delivered', 'in_current_pi', 'planned', 'no_match'):
        effective_status = req.force_status
        member_details = None

    result = await _generate_email(
        req.account_name or primary.get("account_name", "Customer"),
        primary.get("case_number", "N/A"),
        cluster.get("cluster_title", primary.get("subject", "Feature Request")),
        len(member_ids),
        effective_status,
        member_details=member_details,
    )
    result["effective_status"] = effective_status
    result["rfe_subject"] = primary.get("subject", "")
    return result

@app.post("/api/send-emails")
async def api_send_emails():
    conn = get_db(); cur = conn.cursor()
    cur.execute("SELECT * FROM rfe_clusters WHERE status='approved'")
    clusters = [dict(r) for r in cur.fetchall()]
    if not clusters: conn.close(); return {"sent": 0, "message": "No approved clusters"}
    sent_total = 0
    for cluster in clusters:
        run_id = cluster["run_id"]
        member_ids = json.loads(cluster.get("member_case_numbers") or "[]")
        if not member_ids: continue
        ph = ",".join("?" * len(member_ids))
        cur.execute(f"SELECT id,account_name,case_number,subject FROM rfe_pulls WHERE id IN ({ph})", member_ids)
        members = [dict(r) for r in cur.fetchall()]
        output_dir = Path(f"output/emails_{run_id}"); output_dir.mkdir(parents=True, exist_ok=True)
        # Build per-member status lookup from score_breakdown
        try:
            sb = json.loads(cluster.get("score_breakdown") or "{}")
            mm_by_id = {mm["id"]: mm for mm in sb.get("member_matches", [])}
        except Exception:
            mm_by_id = {}
        # Group members by account so each account gets one combined email
        from collections import defaultdict as _dd
        by_account: Dict[str, list] = _dd(list)
        for m in members:
            by_account[m.get("account_name") or "Unknown"].append(m)
        _priority = ["delivered", "in_current_pi", "planned", "no_match"]
        seen = set()
        for account, acct_members in by_account.items():
            if account in seen: continue
            seen.add(account)
            # Build per-RFE detail list for this account
            member_details = []
            for m in acct_members:
                mm = mm_by_id.get(m["id"])
                st = (mm.get("match_status") if mm and mm.get("match_status") else None) \
                     or cluster.get("match_status", "no_match")
                member_details.append({"case_number": m.get("case_number","N/A"),
                                       "subject": m.get("subject",""), "match_status": st})
            effective_status = min(
                (d["match_status"] for d in member_details),
                key=lambda s: _priority.index(s) if s in _priority else 99,
                default=cluster.get("match_status", "no_match")
            )
            primary_m = acct_members[0]
            try:
                email = await _generate_email(account, primary_m.get("case_number", "N/A"),
                    cluster.get("cluster_title", primary_m.get("subject", "")), len(member_ids),
                    effective_status, member_details=member_details)
                fname = f"{account.replace('/', '_').replace(' ', '_')}_{member.get('case_number','x')}.txt"
                with open(output_dir / fname, "w") as f:
                    f.write(f"To: {email.get('to','')}\nSubject: {email.get('subject','')}\n\n{email.get('body','')}\n")
                cur.execute(
                    "INSERT INTO email_log (run_id,cluster_id,case_number,account_name,recipient_email,subject,body,status,sent_at) "
                    "VALUES (?,?,?,?,?,?,?,'sent',?)",
                    (run_id, cluster["id"], member.get("case_number"), account,
                     email.get("to",""), email.get("subject",""), email.get("body",""), datetime.utcnow().isoformat()),
                )
                sent_total += 1
            except: pass
        cur.execute("UPDATE rfe_clusters SET status='sent' WHERE id=?", (cluster["id"],))
        cur.execute("UPDATE run_meta SET emails_sent=emails_sent+? WHERE run_id=?", (len(seen), run_id))
    conn.commit(); conn.close()
    return {"sent": sent_total, "message": f"Wrote {sent_total} email files to output/"}

# ─── Stats / history / config ─────────────────────────────────────────────────

@app.get("/api/stats")
async def api_stats():
    conn = get_db(); cur = conn.cursor()
    cur.execute("SELECT run_id FROM run_meta ORDER BY started_at DESC LIMIT 1")
    row = cur.fetchone()
    if not row: conn.close(); return {"total_clusters": 0, "high_confidence": 0, "approved": 0, "total_arr": 0, "run_id": None}
    run_id = row["run_id"]
    cur.execute("""SELECT
                   SUM(CASE WHEN match_status!='no_match' THEN 1 ELSE 0 END) as total,
                   SUM(CASE WHEN confidence_band='high' AND match_status!='no_match' THEN 1 ELSE 0 END) as high,
                   COALESCE(SUM(CASE WHEN match_status!='no_match' THEN total_arr ELSE 0 END),0) as arr,
                   SUM(CASE WHEN match_status='delivered' THEN 1 ELSE 0 END) as delivered,
                   SUM(CASE WHEN match_status='in_current_pi' THEN 1 ELSE 0 END) as in_current_pi,
                   SUM(CASE WHEN match_status='planned' THEN 1 ELSE 0 END) as planned,
                   SUM(CASE WHEN match_status='no_match' THEN 1 ELSE 0 END) as no_match_clusters,
                   COALESCE(SUM(CASE WHEN match_status!='no_match' THEN json_array_length(member_case_numbers) ELSE 0 END),0) as matched_rfe_count
                   FROM rfe_clusters WHERE run_id=?""", (run_id,))
    s = dict(cur.fetchone())
    cur.execute("SELECT rfe_count,status,source,filtered_rfe_count,started_at,completed_at FROM run_meta WHERE run_id=?", (run_id,))
    meta = dict(cur.fetchone() or {})
    conn.close()
    return {**s, "run_id": run_id,
            "rfe_count": meta.get("rfe_count", 0),
            "filtered_rfe_count": meta.get("filtered_rfe_count", 0),
            "run_status": meta.get("status", ""),
            "source": meta.get("source", ""),
            "started_at": meta.get("started_at", ""),
            "completed_at": meta.get("completed_at", ""),
            "anthropic_key_set": bool(os.getenv("ANTHROPIC_API_KEY", "").strip())}

@app.get("/api/email-queue")
async def api_email_queue():
    conn = get_db(); cur = conn.cursor()
    cur.execute("SELECT c.*,r.started_at as pull_date FROM rfe_clusters c JOIN run_meta r ON c.run_id=r.run_id "
                "WHERE c.status='approved' ORDER BY c.created_at DESC")
    clusters = []
    for row in cur.fetchall():
        c = dict(row); c["case_count"] = len(json.loads(c.get("member_case_numbers") or "[]"))
        clusters.append(c)
    conn.close(); return {"queue": clusters}

@app.get("/api/runs")
async def api_runs():
    conn = get_db()
    runs = [dict(r) for r in conn.execute("SELECT * FROM run_meta ORDER BY started_at DESC LIMIT 50").fetchall()]
    conn.close(); return {"runs": runs}

@app.get("/api/config")
async def api_get_config():
    return {
        "similarity_threshold": float(get_cfg("SIMILARITY_THRESHOLD", "0.55")),
        "high_confidence_threshold": float(get_cfg("HIGH_CONFIDENCE_THRESHOLD", "0.85")),
        "pull_days_back": int(get_cfg("PULL_DAYS_BACK", "90")),
        "sf_username": os.getenv("SF_USERNAME", ""),
        "anthropic_key_set": bool(os.getenv("ANTHROPIC_API_KEY", "").strip()),
        "confluence_configured": bool(os.getenv("CONFLUENCE_BASE_URL", "").strip()),
    }

class ConfigUpdate(BaseModel):
    similarity_threshold: Optional[float] = None
    high_confidence_threshold: Optional[float] = None
    pull_days_back: Optional[int] = None

@app.put("/api/config")
async def api_set_config(cfg: ConfigUpdate):
    if cfg.similarity_threshold is not None: set_cfg("SIMILARITY_THRESHOLD", str(cfg.similarity_threshold))
    if cfg.high_confidence_threshold is not None: set_cfg("HIGH_CONFIDENCE_THRESHOLD", str(cfg.high_confidence_threshold))
    if cfg.pull_days_back is not None: set_cfg("PULL_DAYS_BACK", str(cfg.pull_days_back))
    return {"saved": True}

# ─── Legacy SF pull (phase 2) ─────────────────────────────────────────────────

@app.post("/api/pull")
async def api_pull(background_tasks: BackgroundTasks, days_back: Optional[int] = None):
    try: from simple_salesforce import Salesforce
    except ImportError: raise HTTPException(501, "simple-salesforce not installed")
    sf_user = os.getenv("SF_USERNAME","").strip(); sf_pass = os.getenv("SF_PASSWORD","").strip()
    sf_token = os.getenv("SF_SECURITY_TOKEN","").strip(); sf_domain = os.getenv("SF_DOMAIN","login").strip()
    if not sf_user or not sf_pass: raise HTTPException(400, "Salesforce credentials not configured")
    run_id = str(uuid.uuid4()); job_id = str(uuid.uuid4())
    if days_back is None: days_back = int(get_cfg("PULL_DAYS_BACK","90"))
    jobs[job_id] = {"status":"starting","message":"Connecting to Salesforce…","run_id":run_id,"job_id":job_id}
    conn = get_db()
    conn.execute("INSERT OR REPLACE INTO run_meta (run_id,started_at,status) VALUES (?,?,'running')",
                 (run_id, datetime.utcnow().isoformat()))
    conn.commit(); conn.close()
    async def _run():
        loop = asyncio.get_event_loop()
        try:
            await loop.run_in_executor(executor, _sf_pull_sync, run_id, days_back, job_id, sf_user, sf_pass, sf_token, sf_domain)
            await score_run(run_id, job_id)
        except Exception as e:
            jobs[job_id].update(status="error", message=str(e))
    background_tasks.add_task(_run)
    return {"job_id": job_id, "run_id": run_id}

def _sf_pull_sync(run_id, days_back, job_id, sf_user, sf_pass, sf_token, sf_domain):
    from simple_salesforce import Salesforce
    jobs[job_id].update(status="pulling", message="Running SOQL query…")
    sf = Salesforce(username=sf_user, password=sf_pass, security_token=sf_token, domain=sf_domain)
    soql = f"""SELECT CaseNumber,Subject,Description,Account.Name,Account.AnnualRevenue,
               Status,Product_Domain__c,Sub_Domain__c,CreatedDate,Severity__c
               FROM Case WHERE RecordType.Name='RFE'
               AND Status IN ('Opened','Under PM Review')
               AND CreatedDate >= LAST_N_DAYS:{days_back}"""
    records = sf.query_all(soql).get("records", [])
    pulled_at = datetime.utcnow().isoformat()
    conn = get_db(); cur = conn.cursor()
    for rec in records:
        acct = rec.get("Account") or {}
        cur.execute(
            "INSERT INTO rfe_pulls (run_id,case_number,subject,description,account_name,account_arr,"
            "status,domain,sub_domain,severity,created_date,pulled_at) VALUES (?,?,?,?,?,?,?,?,?,?,?,?)",
            (run_id, rec.get("CaseNumber"), rec.get("Subject"), rec.get("Description"),
             acct.get("Name"), acct.get("AnnualRevenue"), rec.get("Status"),
             rec.get("Product_Domain__c"), rec.get("Sub_Domain__c"),
             rec.get("Severity__c"), rec.get("CreatedDate"), pulled_at),
        )
    cur.execute("INSERT OR REPLACE INTO run_meta (run_id,started_at,rfe_count,status,source) VALUES (?,?,?,'pulled','salesforce')",
                (run_id, pulled_at, len(records)))
    conn.commit(); conn.close()
    jobs[job_id].update(status="pulled", message=f"Pulled {len(records)} RFEs", run_id=run_id, rfe_count=len(records))

# ─── Serve frontend ───────────────────────────────────────────────────────────

app.mount("/static", StaticFiles(directory="static"), name="static")

# ── Server-side cluster HTML rendering (mirrors JS buildSignalRow / renderSignalTable)
def _build_clusters_for_page(conn, run_id: str,  # noqa: C901
                              search_query: str = ""):
    """
    Returns (html_str, clusters_list) where:
    - html_str is the pre-rendered HTML for signal-table-container
    - clusters_list is the JSON-serialisable list to inject as window.__INITIAL_CLUSTERS__
    search_query: if set, filter clusters to those matching case ID / subject / keyword
    """
    import html as _h
    E = _h.escape          # same as JS esc()
    import re as _re

    cur = conn.cursor()

    # ── Fetch matched clusters (same logic as /api/clusters) ──────────────────
    cur.execute(
        "SELECT * FROM rfe_clusters WHERE run_id=? AND match_status!='no_match' "
        "ORDER BY confidence_score DESC", (run_id,)
    )
    clusters = [dict(r) for r in cur.fetchall()]
    if not clusters:
        return '', []

    for c in clusters:
        member_ids = json.loads(c.get("member_case_numbers") or "[]")
        if member_ids:
            ph = ",".join("?" * len(member_ids))
            cur.execute(
                f"SELECT id,case_number,subject,description,account_name,"
                f"account_arr,created_date,status,domain "
                f"FROM rfe_pulls WHERE id IN ({ph})", member_ids
            )
            c["members"] = [dict(r) for r in cur.fetchall()]
        else:
            c["members"] = []

        try:
            c["score_breakdown"] = json.loads(c.get("score_breakdown") or "{}")
        except Exception:
            c["score_breakdown"] = {}

        # Attach per-RFE match data to each member (mirrors api_clusters logic)
        mm_by_id = {mm["id"]: mm for mm in c["score_breakdown"].get("member_matches", [])}
        for m in c["members"]:
            mm = mm_by_id.get(m["id"])
            if mm:
                m["match_status"]  = mm.get("match_status", "no_match")
                m["ctx_breakdown"] = mm.get("ctx_breakdown", {})
            else:
                m["match_status"]  = c.get("match_status", "no_match")
                m["ctx_breakdown"] = {}

        if not c.get("confidence_sentence"):
            c["confidence_sentence"] = generate_confidence_sentence(
                c.get("cluster_title", ""), c.get("match_status", "no_match"),
                c.get("match_reason", ""), c.get("confidence_score", 0)
            )

    # ── Apply search filter ───────────────────────────────────────────────────
    if search_query:
        q = search_query.lower().strip()
        filtered = []
        for c in clusters:
            # Match on cluster title
            if q in (c.get("cluster_title") or "").lower():
                filtered.append(c); continue
            # Match on any member case number, subject, or description
            matched = False
            for m in c.get("members", []):
                if (q in (m.get("case_number") or "").lower() or
                        q in (m.get("subject") or "").lower() or
                        q in (m.get("description") or "").lower() or
                        q in (m.get("account_name") or "").lower()):
                    matched = True; break
            if matched:
                filtered.append(c)
        clusters = filtered

    if not clusters:
        if search_query:
            return (
                '<div class="empty-state"><div class="big-icon">🔎</div>'
                '<h3>No results</h3>'
                f'<p>No clusters matched <strong>{_h.escape(search_query)}</strong>'
                ' — try a different case ID, subject, or keyword.</p></div>'
            ), []
        return '', []

    # ── Map each cluster to a canonical domain & annotate ────────────────────
    for c in clusters:
        raw_dom = (c["members"][0].get("domain") if c["members"] else None) or ""
        c["_canonical_domain"] = _map_domain(raw_dom, c.get("cluster_title", ""))

    # ── Group by canonical domain (preserve CANONICAL_DOMAINS order) ─────────
    groups: dict = {}
    for c in clusters:
        grp = c["_canonical_domain"]
        groups.setdefault(grp, []).append(c)
    # Sort groups by the canonical order; unknowns go last
    def _grp_sort_key(k):
        try: return CANONICAL_DOMAINS.index(k)
        except ValueError: return 999
    groups = dict(sorted(groups.items(), key=lambda kv: _grp_sort_key(kv[0])))

    _STATUS_LABELS  = {'delivered': '✅ Delivered', 'in_current_pi': '🏃 In PI', 'planned': '🗓 Planned'}
    _STATUS_CLASSES = {'delivered': 'status-delivered', 'in_current_pi': 'status-in_current_pi', 'planned': 'status-planned'}
    _CONF_DOTS      = {'high': '🟢', 'medium': '🟡', 'low': '🔴'}
    _RFE_GRID       = '72px 1fr 56px 105px 130px 145px'

    # Source badge styles (matches JS buildMatchedSources)
    _SRC_STYLE = {
        'release_notes':  'background:#f3e8ff;color:#7e22ce;border:1px solid #d8b4fe',
        'ado_current_pi': 'background:#dbeafe;color:#1d4ed8;border:1px solid #93c5fd',
        'ado_backlog':    'background:#dbeafe;color:#1d4ed8;border:1px solid #93c5fd',
        'confluence':     'background:#ccfbf1;color:#0f766e;border:1px solid #99f6e4',
        'roadmap':        'background:#ffedd5;color:#c2410c;border:1px solid #fdba74',
    }
    _SRC_LABEL = {
        'release_notes':  'RELEASE NOTES',
        'ado_current_pi': 'ADO EPIC',
        'ado_backlog':    'ADO EPIC',
        'confluence':     'CONFLUENCE',
        'roadmap':        'ROADMAP',
    }
    _MS_STATUS_PFX = {
        'delivered':     ['release_notes'],
        'in_current_pi': ['ado_current_pi'],
        'planned':       ['ado_backlog', 'roadmap', 'confluence'],
    }
    _MS_MARK = {'delivered': '✔ Delivered', 'in_current_pi': '✔ In PI', 'planned': '✔ Planned'}

    def _arr_fmt(v):
        v = v or 0
        if v >= 1e6: return f'${v/1e6:.1f}M'
        if v >= 1e3: return f'${v/1e3:.0f}K'
        if v > 0:    return f'${v:.0f}'
        return '—'

    def _matched_sources_html(m_ctx_bd, cluster_bd, m_ms):
        """Render matched-source pills for one RFE (mirrors JS buildMatchedSources)."""
        bd = m_ctx_bd if m_ctx_bd else (cluster_bd or {})
        has_per_rfe = bool(m_ctx_bd)

        # Sort by score desc, keep one per docType
        sorted_entries = sorted(bd.items(), key=lambda x: -x[1])
        seen_types: set = set()
        entries = []
        for key, score in sorted_entries:
            if score < 0.05:
                continue
            sep = key.find('|')
            dt = key[:sep] if sep > -1 else key
            if dt in seen_types:
                continue
            seen_types.add(dt)
            entries.append((key, score, dt))

        if not entries:
            return ('<div style="margin-top:6px;font-size:11px;color:var(--text-muted);'
                    'font-style:italic">No specific source match found.</div>')

        status_pfx = _MS_STATUS_PFX.get(m_ms, [])
        ms_mark    = _MS_MARK.get(m_ms, '')
        pills = []
        for key, score, dt in entries:
            sep   = key.find('|')
            title = key[sep+1:] if sep > -1 else key
            b_sty = _SRC_STYLE.get(dt, 'background:#f1f5f9;color:var(--text-muted);border:1px solid var(--border)')
            b_lbl = _SRC_LABEL.get(dt, dt.upper().replace('_', ' '))
            is_src = any(key.startswith(p) for p in status_pfx)

            # ADO epics: show "Epic Name #ID" with #ID copyable
            if dt in ('ado_current_pi', 'ado_backlog'):
                hash_m = _re.search(r' (#\d+)$', title)
                if hash_m:
                    epic_name = E(title[:title.rfind(hash_m.group(0))])
                    epic_id   = E(hash_m.group(1))
                    title_html = (
                        epic_name +
                        ' <span data-cpid="' + epic_id + '" style="cursor:pointer;color:var(--accent);font-weight:700"'
                        ' onclick="event.stopPropagation();copyMatchedId(this,this.getAttribute(\'data-cpid\'))"'
                        ' title="Click to copy">' + epic_id + '</span>'
                    )
                else:
                    safe_t = E(title)
                    title_html = (
                        '<span data-cpid="' + safe_t + '" style="cursor:pointer;color:var(--text)"'
                        ' onclick="event.stopPropagation();copyMatchedId(this,this.getAttribute(\'data-cpid\'))"'
                        ' title="Click to copy">' + safe_t + '</span>'
                    )
            else:
                safe_t = E(title)
                title_html = (
                    '<span data-cpid="' + safe_t + '" style="cursor:pointer;color:var(--text)"'
                    ' onclick="event.stopPropagation();copyMatchedId(this,this.getAttribute(\'data-cpid\'))"'
                    ' title="Click to copy">' + safe_t + '</span>'
                )

            status_badge = ''
            if is_src and ms_mark:
                status_badge = (
                    '<span style="background:#dcfce7;color:#15803d;border:1px solid #86efac;'
                    'border-radius:3px;padding:1px 5px;font-size:9px;font-weight:700">'
                    + ms_mark + '</span>'
                )
            pills.append(
                '<div style="display:inline-flex;align-items:center;gap:5px;background:#f8fafc;'
                'border:1px solid var(--border);border-radius:6px;padding:3px 8px;font-size:11px">'
                '<span style="' + b_sty + ';border-radius:3px;padding:1px 5px;font-size:9px;font-weight:700;white-space:nowrap">'
                + b_lbl + '</span>'
                + title_html
                + '<span style="color:var(--text-muted)">[' + f'{score:.2f}' + ']</span>'
                + status_badge
                + '</div>'
            )

        src_label = ('Matched Sources'
                     if has_per_rfe else
                     'Matched Sources <span style="font-size:9px;font-weight:400;color:var(--text-muted);font-style:italic">(cluster-level)</span>')
        return (
            '<div style="margin-top:6px">'
            '<div style="font-size:10px;font-weight:700;color:var(--text-muted);margin-bottom:4px;'
            'text-transform:uppercase;letter-spacing:.05em">' + src_label + '</div>'
            '<div style="display:flex;flex-wrap:wrap;gap:5px">' + ''.join(pills) + '</div>'
            '</div>'
        )

    def _build_row(c):
        members   = c.get("members", [])
        accounts  = list(dict.fromkeys(m["account_name"] for m in members if m.get("account_name")))
        domain    = c.get("_canonical_domain") or (members[0].get("domain") if members else None) or '—'
        ms        = c.get("match_status", "no_match")
        conf_band = c.get("confidence_band", "low")
        conf_scr  = c.get("confidence_score", 0) or 0
        cid       = c["id"]

        conf_chip = (f'<span class="chip chip-{E(conf_band)}">'
                     f'{_CONF_DOTS.get(conf_band, "")} {conf_scr:.2f}</span>')

        # ── Confidence sentence ─────────────────────────────────────────────
        sent = c.get("confidence_sentence", "") or ""
        sentence_html = (
            f'<div style="margin-bottom:14px;font-style:italic;color:var(--text-muted);'
            f'font-size:13px;line-height:1.5;border-left:3px solid var(--accent);'
            f'padding-left:10px">{E(sent)}</div>'
        ) if sent else ''

        # ── RFE grid ─────────────────────────────────────────────────────────
        ctx_bd = c.get("score_breakdown", {}).get("context_breakdown", {})
        rfe_header = (
            f'<div style="display:grid;grid-template-columns:{_RFE_GRID};gap:8px;'
            f'font-size:10px;font-weight:700;color:var(--text-muted);text-transform:uppercase;'
            f'letter-spacing:.05em;padding:6px 0;border-bottom:1px solid var(--border);margin-bottom:2px">'
            f'<div style="text-align:center">Case #</div><div>Subject</div>'
            f'<div style="text-align:center">Signal</div>'
            f'<div style="text-align:center">Match Status</div>'
            f'<div style="text-align:center">RFE State</div>'
            f'<div style="text-align:center">Customer</div></div>'
        )

        rfe_rows = rfe_header
        for m in members:
            case_id  = m.get("case_number") or "—"
            m_ms     = m.get("match_status") or ms
            m_chip   = (
                '<span class="status-chip ' + _STATUS_CLASSES.get(m_ms, '') + '">'
                + _STATUS_LABELS.get(m_ms, m_ms) + '</span>'
            ) if m_ms != "no_match" else "—"

            _m_bd      = m.get("ctx_breakdown") or ctx_bd or {}
            _m_entries = sorted(_m_bd.items(), key=lambda x: -x[1]) if _m_bd else []
            _m_signal  = f"{_m_entries[0][1]:.2f}" if _m_entries else "—"

            desc = (m.get("description") or "").strip()
            desc_html = (
                '<div>' + E(desc[:400]) + '</div>' if desc else
                '<div style="color:var(--text-muted);font-style:italic">No description provided.</div>'
            )

            safe_case   = E(case_id).replace("'", "&#39;")
            m_subj_e    = E(m.get("subject", "") or "")
            m_subj_disp = E(m.get("subject", "") or "—")
            m_status_e  = E(m.get("status") or "—")
            m_acct_e    = E(m.get("account_name") or "—")

            # Matched sources pills for this RFE
            src_html = _matched_sources_html(
                m.get("ctx_breakdown") or {},
                ctx_bd,
                m_ms
            )

            rfe_rows += (
                '<div style="background:#f4f6f9;border-radius:6px;padding:8px 10px;margin-bottom:4px">'
                '<div style="display:grid;grid-template-columns:' + _RFE_GRID + ';gap:8px;align-items:center;font-size:12px">'
                '<span class="copy-case" onclick="copyCase(this,\'' + safe_case + '\')" title="Click to copy" '
                'style="display:flex;align-items:center;justify-content:center;text-align:center">' + E(case_id) + '</span>'
                '<span style="font-weight:600;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;min-width:0" '
                'title="' + m_subj_e + '">' + m_subj_disp + '</span>'
                '<span style="color:var(--accent);font-weight:700;text-align:center;display:block">' + _m_signal + '</span>'
                '<div style="display:flex;justify-content:center">' + m_chip + '</div>'
                '<div style="overflow:hidden;text-overflow:ellipsis;white-space:nowrap;color:var(--text-muted);text-align:center">'
                + m_status_e + '</div>'
                '<div style="overflow:hidden;text-overflow:ellipsis;white-space:nowrap;color:var(--text-muted);text-align:center">'
                + m_acct_e + '</div>'
                '</div>'
                '<div style="margin-top:6px;padding-left:2px;font-size:12px;line-height:1.5">'
                '<div style="display:flex;align-items:center;gap:10px;margin-bottom:4px;flex-wrap:wrap">'
                '<strong style="font-weight:700;color:var(--text)">Summary:</strong></div>'
                + desc_html
                + src_html
                + '</div></div>'
            )

        # ── Email section: customer dropdown (multi-account) ──────────────────
        # json.dumps produces "Name" with double-quotes; &quot; makes it safe
        # inside an HTML onclick="…" attribute.
        def _safe_js_str(s):
            """JSON-encode a string and make it safe inside an HTML attribute."""
            return json.dumps(s).replace('"', '&quot;')

        if len(accounts) <= 1:
            acct_arg = _safe_js_str(accounts[0]) if accounts else 'null'
            email_html = (
                f'<button class="btn btn-primary btn-sm" '
                f'onclick="event.stopPropagation();openEmailDrawer({cid},{acct_arg})">'
                f'✉ Generate Email</button>'
            )
        else:
            opts   = "".join(f'<option value="{E(a)}">{E(a)}</option>' for a in accounts)
            sel_id = f'acct-sel-{cid}'
            sel_get = "document.getElementById('" + sel_id + "').value"
            email_html = (
                '<div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap">'
                '<label style="font-size:11px;font-weight:700;color:var(--text-muted)">Customer:</label>'
                '<select id="' + sel_id + '" onclick="event.stopPropagation()" '
                'style="padding:5px 8px;border:1px solid var(--border);border-radius:7px;font-size:12px;background:var(--card)">'
                + opts + '</select>'
                '<button class="btn btn-primary btn-sm" '
                'onclick="event.stopPropagation();openEmailDrawer(' + str(cid) + ',' + sel_get + ')">'
                '✉ Generate Email</button></div>'
            )

        # ── Customer display in summary row ───────────────────────────────────
        if accounts:
            cust_html = E(accounts[0])
            if len(accounts) > 1:
                cust_html += f'<br><span style="font-size:10px;color:var(--text-muted)">+{len(accounts)-1} more</span>'
        else:
            cust_html = '—'

        mt = '12px' if sentence_html else '0'
        return (
            f'\n<div class="signal-row" id="row-{cid}" onclick="toggleExpand({cid})">'
            f'<div><span class="sig-domain">{E(domain)}</span></div>'
            f'<div><div class="sig-title">{E(c.get("cluster_title") or "Untitled")}</div>'
            f'<div class="sig-accounts">{len(members)} RFE{"s" if len(members)!=1 else ""}</div></div>'
            f'<div style="font-size:12px;font-weight:500;overflow:hidden">{cust_html}</div>'
            f'<div style="text-align:center">{conf_chip}</div>'
            f'<div class="sig-arr">{_arr_fmt(c.get("total_arr"))}</div>'
            f'</div>'
            f'<div class="why-row" id="expand-{cid}" style="display:none">'
            f'<div class="why-panel">'
            f'{sentence_html}'
            f'<div style="margin-bottom:14px;margin-top:{mt}">'
            f'<div class="why-label" style="margin-bottom:6px">📋 RFEs in this cluster</div>'
            f'{rfe_rows}</div>'
            f'<div style="padding-top:10px;border-top:1px solid var(--border)">'
            f'{email_html}</div>'
            f'</div></div>'
        )

    # ── Assemble HTML (groups already sorted by canonical order) ─────────────
    html_parts = []
    for grp_key in groups.keys():
        grp_clusters = groups[grp_key]
        n = len(grp_clusters)
        rows_html = "".join(_build_row(c) for c in grp_clusters)
        html_parts.append(
            f'<div class="signal-table-wrap" style="margin-bottom:12px">'
            f'<div class="signal-group-header">'
            f'<span>{E(grp_key)}</span>'
            f'<span style="font-weight:400">{n} cluster{"s" if n!=1 else ""}</span>'
            f'</div>'
            f'<div>'
            f'<div class="signal-row" style="background:#fafafa;cursor:default">'
            f'<div style="font-size:10px;font-weight:700;color:var(--text-muted);text-transform:uppercase">Domain</div>'
            f'<div style="font-size:10px;font-weight:700;color:var(--text-muted);text-transform:uppercase">Cluster / RFEs</div>'
            f'<div style="font-size:10px;font-weight:700;color:var(--text-muted);text-transform:uppercase">Customer</div>'
            f'<div style="font-size:10px;font-weight:700;color:var(--text-muted);text-transform:uppercase;text-align:center">Score</div>'
            f'<div style="font-size:10px;font-weight:700;color:var(--text-muted);text-transform:uppercase">ARR</div>'
            f'</div>'
            f'{rows_html}'
            f'</div></div>'
        )

    return "".join(html_parts), clusters

@app.get("/")
async def root():
    """Serve the app with all state server-side rendered — no JS fetch required for initial display."""
    import html as _esc

    # ── 1. Load DB state ────────────────────────────────────────────────────────
    stats = {}  # Signal Match stat card values
    try:
        conn = get_db()
        counts = {r["doc_type"]: r["cnt"] for r in conn.execute(
            "SELECT doc_type, COUNT(*) as cnt FROM context_docs GROUP BY doc_type"
        ).fetchall()}
        docs = [dict(r) for r in conn.execute(
            "SELECT id, doc_type, title, imported_at, length(content) as chars "
            "FROM context_docs ORDER BY imported_at DESC"
        ).fetchall()]
        run_row = conn.execute(
            "SELECT run_id, rfe_count, status, filtered_rfe_count FROM run_meta "
            "ORDER BY started_at DESC LIMIT 1"
        ).fetchone()
        run_meta = dict(run_row) if run_row else {}
        if run_meta.get("run_id"):
            cl_row = conn.execute(
                "SELECT COUNT(*) as cnt FROM rfe_clusters WHERE run_id=?",
                (run_meta["run_id"],)
            ).fetchone()
            run_meta["cluster_count"] = cl_row["cnt"] if cl_row else 0
            # Signal Match stats
            st_row = conn.execute("""
                SELECT
                  SUM(CASE WHEN match_status!='no_match' THEN 1 ELSE 0 END) as total,
                  SUM(CASE WHEN confidence_band='high' AND match_status!='no_match' THEN 1 ELSE 0 END) as high,
                  COALESCE(SUM(CASE WHEN match_status!='no_match' THEN total_arr ELSE 0 END),0) as arr,
                  SUM(CASE WHEN match_status='delivered' THEN 1 ELSE 0 END) as delivered,
                  SUM(CASE WHEN match_status='in_current_pi' THEN 1 ELSE 0 END) as in_current_pi,
                  COALESCE(SUM(CASE WHEN match_status!='no_match'
                    THEN json_array_length(member_case_numbers) ELSE 0 END),0) as matched_rfe_count
                FROM rfe_clusters WHERE run_id=?""", (run_meta["run_id"],)
            ).fetchone()
            if st_row:
                stats = dict(st_row)
                stats["rfe_count"] = run_meta.get("rfe_count", 0)
                arr = stats.get("arr", 0) or 0
                stats["arr_fmt"] = (
                    f'${arr/1e6:.1f}M' if arr >= 1e6 else
                    f'${arr/1e3:.0f}K' if arr >= 1e3 else
                    f'${arr:.0f}' if arr > 0 else '—'
                )
        conn.close()
    except Exception as _e:
        counts = {}; docs = []; run_meta = {}; stats = {}

    page = Path("static/index.html").read_text(encoding="utf-8")

    # ── 2. JS bootstrap variables (injected before any script runs) ─────────────
    run_id_val = f'"{run_meta.get("run_id","")}"' if run_meta.get("run_id") else 'null'
    page = page.replace("</head>",
        f"<script>"
        f"window.__INITIAL_COUNTS__={json.dumps(counts)};"
        f"window.__INITIAL_RUN_ID__={run_id_val};"
        f"</script>\n</head>", 1)

    # ── 3. Source-card status badges (simple string replacements) ───────────────
    rn    = counts.get("release_notes", 0)
    rm    = counts.get("roadmap", 0)
    cf    = counts.get("confluence", 0)
    adopi = counts.get("ado_current_pi", 0)
    adobl = counts.get("ado_backlog", 0)
    total = sum(counts.values())

    def _status(page, sid, cid, text):
        page = re.sub(
            r'<div class="source-status[^"]*" id="' + sid + r'">[^<]*</div>',
            f'<div class="source-status loaded" id="{sid}">{text}</div>',
            page, count=1
        )
        page = page.replace(f'id="{cid}" class="source-card"',
                            f'id="{cid}" class="source-card has-data"', 1)
        page = page.replace(f'class="source-card" id="{cid}"',
                            f'class="source-card has-data" id="{cid}"', 1)
        return page

    if rn    > 0: page = _status(page,"status-release_notes","card-release_notes",
                                  f"✅ {rn} version{'s'if rn!=1 else ''} loaded")
    if rm    > 0: page = _status(page,"status-roadmap","card-roadmap",
                                  f"✅ {rm} file{'s'if rm!=1 else''} loaded")
    if cf    > 0: page = _status(page,"status-confluence","card-confluence",
                                  f"✅ {cf} page{'s'if cf!=1 else''} loaded")
    if adopi > 0 or adobl > 0:
        parts=[]
        if adopi>0: parts.append(f"Current PI: {adopi}")
        if adobl >0: parts.append(f"Upcoming PI: {adobl}")
        page = _status(page,"status-ado","card-ado","✅ "+" · ".join(parts))

    # ── 4. Nav badges ────────────────────────────────────────────────────────────
    if total > 0:
        page = page.replace(
            '<span class="nav-badge" id="badge-sources" style="display:none"></span>',
            f'<span class="nav-badge" id="badge-sources">{total}</span>', 1)
    # Signal Match badge = matched (non-no_match) cluster count
    sig_badge_n = stats.get("total", 0) or 0
    if sig_badge_n > 0:
        page = page.replace(
            '<span class="nav-badge" id="badge-signal" style="display:none">0</span>',
            f'<span class="nav-badge" id="badge-signal">{sig_badge_n}</span>', 1)

    # ── 5. Loaded-docs section ───────────────────────────────────────────────────
    if docs:
        ICONS  = {"release_notes":"📋","ado_current_pi":"🏃","ado_backlog":"📋",
                  "roadmap":"🗺️","confluence":"📝"}
        LABELS = {"release_notes":"Release Notes","ado_current_pi":"ADO Current PI",
                  "ado_backlog":"ADO Backlog","roadmap":"Roadmap","confluence":"Confluence"}
        rows_html = "".join(
            f'<div class="doc-row">'
            f'<span class="doc-type-badge badge-{d["doc_type"]}">'
            f'{ICONS.get(d["doc_type"],"📄")} {LABELS.get(d["doc_type"],d["doc_type"])}</span>'
            f'<span class="doc-title">{_esc.escape(d["title"] or "")}</span>'
            f'<span class="doc-meta">'
            f'{str(round(d["chars"]/1000))+"k chars · " if d.get("chars") else ""}'
            f'{(d.get("imported_at") or "")[:16].replace("T"," ")}</span>'
            f'<button class="btn btn-ghost btn-sm" onclick="deleteDoc({d["id"]})">✕</button>'
            f'</div>'
            for d in docs
        )
        page = page.replace('id="docs-section" style="display:none"',
                            'id="docs-section"', 1)
        # The docs-list div has both class and id: <div class="docs-list" id="docs-list">
        page = page.replace('<div class="docs-list" id="docs-list"></div>',
                            f'<div class="docs-list" id="docs-list">{rows_html}</div>', 1)

    # ── 6. Signal Match stat cards (pre-populate so they show instantly) ──────────
    if stats:
        def _stat(page, eid, val):
            return page.replace(f'id="{eid}">—<', f'id="{eid}">{val}<', 1)
        page = _stat(page, "stat-total",        stats.get("total") or "—")
        page = _stat(page, "stat-matched-rfes", stats.get("matched_rfe_count") or "—")
        page = _stat(page, "stat-rfe-big",      stats.get("rfe_count") or "—")
        page = _stat(page, "stat-high",         stats.get("high") or "—")
        page = _stat(page, "stat-delivered",    stats.get("delivered") or "—")
        page = _stat(page, "stat-in-pi",        stats.get("in_current_pi") or "—")
        page = _stat(page, "stat-arr",          stats.get("arr_fmt") or "—")

    # ── 7. CSV drop-zone status (restore previous session) ──────────────────────
    run_status = run_meta.get("status", "")
    rfe_count  = run_meta.get("rfe_count", 0)
    n_clusters = run_meta.get("cluster_count", 0)
    if run_status == "scored" and rfe_count:
        cluster_txt = f"{n_clusters} matched clusters" if n_clusters else "0 clusters"
        csv_status  = (f'<span style="color:var(--success)">✅ {rfe_count} RFEs loaded'
                       f' · {cluster_txt} · data saved in DB</span>')
        page = page.replace(
            '<div id="csv-upload-status" style="margin-top:10px;font-size:13px"></div>',
            f'<div id="csv-upload-status" style="margin-top:10px;font-size:13px">{csv_status}</div>',
            1)
        page = page.replace(
            'class="drop-zone" id="csv-drop-zone"',
            'class="drop-zone" id="csv-drop-zone" style="border-color:var(--success)"', 1)

    # ── 8. Signal Match cluster list — server-side render ────────────────────────
    # Build HTML and data for allClusters JS variable so all interactive features work
    clusters_html = ''
    clusters_data = []
    if run_meta.get("run_id"):
        try:
            conn2 = get_db()
            clusters_html, clusters_data = _build_clusters_for_page(conn2, run_meta["run_id"])
            conn2.close()
        except Exception as _ce:
            clusters_html = (
                f'<div class="empty-state"><div class="big-icon">⚠️</div>'
                f'<h3>Failed to render clusters</h3><p>{_ce}</p></div>'
            )
    if clusters_html:
        page = page.replace(
            '<div id="signal-table-container"></div>',
            f'<div id="signal-table-container">{clusters_html}</div>', 1
        )

    # ── 9. Inject allClusters data for interactive JS features (email drawer etc.)
    # Rebuild the head-injection to include __INITIAL_CLUSTERS__
    # We already injected run_id+counts; now replace that injection with an extended one
    clusters_json = json.dumps(clusters_data, default=str)
    page = page.replace(
        f"window.__INITIAL_COUNTS__={json.dumps(counts)};"
        f"window.__INITIAL_RUN_ID__={run_id_val};",
        f"window.__INITIAL_COUNTS__={json.dumps(counts)};"
        f"window.__INITIAL_RUN_ID__={run_id_val};"
        f"window.__INITIAL_CLUSTERS__={clusters_json};",
        1
    )

    return HTMLResponse(content=page,
        headers={"Cache-Control":"no-store, no-cache, must-revalidate","Pragma":"no-cache"})

if __name__ == "__main__":
    import uvicorn, webbrowser
    def _open():
        import time; time.sleep(1.5); webbrowser.open("http://localhost:8000")
    threading.Thread(target=_open, daemon=True).start()
    uvicorn.run(app, host="0.0.0.0", port=8000, reload=False)
